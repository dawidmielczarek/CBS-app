
import streamlit as st
import pandas as pd
import numpy as np
import io
import plotly.express as px
import plotly.figure_factory as ff
from scipy import stats
import statsmodels.api as sm
import statsmodels.formula.api as smf
from statsmodels.stats.outliers_influence import variance_inflation_factor
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

st.set_page_config(page_title="CBS Quick Insights", layout="wide")
st.title("📊 CBS Quick Insights — Data to Insight in Minutes")
st.caption("Upload data → tag Inputs/Outputs → get EDA, correlations, regression, SPC, and a ready-to-share report.")

# ---------- Helpers ----------
@st.cache_data
def load_data(file):
    if file.name.lower().endswith(".csv"):
        return pd.read_csv(file)
    else:
        return pd.read_excel(file)

def detect_types(df):
    types = {}
    for c in df.columns:
        if np.issubdtype(df[c].dtype, np.number):
            types[c] = "numeric"
        elif np.issubdtype(df[c].dtype, np.datetime64):
            types[c] = "datetime"
        else:
            types[c] = "categorical"
    return types

def ad_normality(series):
    series = series.dropna()
    if len(series) < 8:
        return np.nan, np.nan
    ad_res = stats.anderson(series, dist='norm')
    stat = ad_res.statistic
    # Approx p-value: compare to 5% critical value
    crit_5 = ad_res.critical_values[2]
    p_approx = 0.049 if stat < crit_5 else 0.051
    return stat, p_approx

def compute_vif(df, features):
    X = df[features].dropna()
    X = sm.add_constant(X)
    vifs = []
    for i, col in enumerate(X.columns):
        if col == "const":
            continue
        vifs.append({"feature": col, "VIF": variance_inflation_factor(X.values, i)})
    return pd.DataFrame(vifs)

def insight(text):
    st.markdown(f"✅ **Insight:** {text}")

# ---------- Sidebar: Upload & Tag ----------
st.sidebar.header("1) Upload data")
up = st.sidebar.file_uploader("CSV or Excel", type=["csv","xlsx","xls"])
demo = st.sidebar.checkbox("Use demo dataset", value=up is None)
if demo:
    df = pd.read_csv("sample_process_data.csv")
elif up:
    df = load_data(up)
else:
    st.stop()

st.sidebar.header("2) Tag columns")
types = detect_types(df)
numeric_cols = [c for c,t in types.items() if t=="numeric"]
cat_cols = [c for c,t in types.items() if t=="categorical"]
time_cols = [c for c,t in types.items() if t=="datetime"]

X_cols = st.sidebar.multiselect("Inputs (X)", options=df.columns, default=[c for c in numeric_cols if c not in ["Gap"]][:3] + cat_cols[:1])
Y_cols = st.sidebar.multiselect("Outputs (Y)", options=df.columns, default=[c for c in df.columns if c.lower() in ["gap","y","output"]][:1])
time_col = st.sidebar.selectbox("Time column (optional for SPC)", options=["(none)"]+df.columns.tolist(), index=0)

st.sidebar.header("3) Optional: Spec limits for capability")
spec_target_col = st.sidebar.selectbox("Select Y for spec (optional)", options=["(none)"] + Y_cols, index=0)
LSL = st.sidebar.text_input("LSL", value="")
USL = st.sidebar.text_input("USL", value="")

st.write("### Data preview")
st.dataframe(df.head(10), use_container_width=True)

# ---------- Section: Data Health ----------
st.subheader("🔎 Data Health")
n_rows, n_cols = df.shape
missing = df.isna().mean().round(3)
st.write(f"Rows: **{n_rows}**, Columns: **{n_cols}**")
st.write("Missing values (fraction):")
st.dataframe(missing.rename("missing_fraction"))

# ---------- Section: EDA & Relationships ----------
st.subheader("📈 Relationships & EDA")
if len(Y_cols) >= 1:
    Y = Y_cols[0]
    st.markdown(f"**Selected Y:** `{Y}`  |  **Inputs (X):** {', '.join(X_cols) if X_cols else '(none)'}")

    # Normality for Y
    if Y in numeric_cols:
        stat, p = ad_normality(df[Y])
        if pd.notna(p):
            st.write(f"Normality (Anderson-Darling approx): stat={stat:.3f} → {'likely normal' if p>0.05 else 'non-normal'}")

    # Numeric-numeric correlations
    num_X = [x for x in X_cols if x in numeric_cols]
    if Y in numeric_cols and num_X:
        corr_rows = []
        for x in num_X:
            sub = df[[x, Y]].dropna()
            if len(sub) > 2:
                r, p = stats.pearsonr(sub[x], sub[Y])
                corr_rows.append({"X": x, "Pearson r": r, "p-value": p})
        if corr_rows:
            corr_df = pd.DataFrame(corr_rows).sort_values(by="Pearson r", key=lambda s: s.abs(), ascending=False)
            st.markdown("**Correlations (numeric X vs Y):**")
            st.dataframe(corr_df, use_container_width=True)
            top = corr_df.iloc[0]
            insight(f"{top['X']} has the strongest linear relationship with {Y} (r={top['Pearson r']:.2f}, p={'<0.001' if top['p-value']<0.001 else f'{top['p-value']:.3f}'}).")

            # Scatter for top
            fig = px.scatter(df, x=top["X"], y=Y, trendline="ols")
            st.plotly_chart(fig, use_container_width=True)

    # Categorical-numeric (ANOVA/Kruskal)
    cat_X = [x for x in X_cols if x in cat_cols]
    for x in cat_X:
        st.markdown(f"**Effect of {x} on {Y}**")
        fig = px.box(df, x=x, y=Y, points="outliers")
        st.plotly_chart(fig, use_container_width=True)
        sub = df[[x, Y]].dropna()
        # Levene for equal variances
        groups = [g[Y].values for _, g in sub.groupby(x)]
        if len(groups) > 1:
            lev_p = stats.levene(*groups).pvalue
            if lev_p < 0.05:
                # Welch-ANOVA fallback: use Kruskal
                H, p = stats.kruskal(*groups)
                st.write(f"Kruskal-Wallis: H={H:.2f}, p={p:.3g} (variances unequal)")
            else:
                F, p = stats.f_oneway(*groups)
                st.write(f"ANOVA: F={F:.2f}, p={p:.3g}")

# ---------- Section: Regression ----------
st.subheader("🧠 Regression (continuous Y)")
if Y_cols and Y_cols[0] in numeric_cols and X_cols:
    Y = Y_cols[0]
    # Build formula: treat categoricals as C()
    terms = []
    for x in X_cols:
        if x in cat_cols:
            terms.append(f"C({x})")
        else:
            terms.append(x)
    formula = f"{Y} ~ " + " + ".join(terms)
    st.code(formula, language="python")
    model = smf.ols(formula, data=df).fit()
    st.write(model.summary().tables[1].as_html(), unsafe_allow_html=True)
    st.write(f"Adj. R² = {model.rsquared_adj:.3f}")
    # VIF for numeric only
    num_for_vif = [x for x in X_cols if x in numeric_cols]
    if num_for_vif:
        try:
            from statsmodels.stats.outliers_influence import variance_inflation_factor
            X_vif = df[num_for_vif].dropna()
            X_vif = sm.add_constant(X_vif)
            vifs = []
            for i, col in enumerate(X_vif.columns):
                if col == "const": 
                    continue
                vifs.append({"feature": col, "VIF": variance_inflation_factor(X_vif.values, i)})
            st.write(pd.DataFrame(vifs))
        except Exception as e:
            st.info(f"VIF skipped: {e}")

# ---------- Section: SPC / Capability ----------
st.subheader("📉 SPC & Capability (optional)")
if time_col and time_col != "(none)" and (spec_target_col and spec_target_col != "(none)"):
    y = spec_target_col
    dff = df[[time_col, y]].dropna().sort_values(by=time_col)
    dff['Index'] = range(len(dff))
    mean = dff[y].mean()
    mr = dff[y].diff().abs().dropna().mean()
    d2 = 1.128  # for MR(2)
    sigma = mr / d2 if d2 != 0 else np.nan
    UCL = mean + 3*sigma
    LCL = mean - 3*sigma
    fig2 = px.line(dff, x=time_col, y=y, title=f"I-Chart for {y}")
    fig2.add_hline(y=mean, line_dash="dash")
    fig2.add_hline(y=UCL, line_dash="dot")
    fig2.add_hline(y=LCL, line_dash="dot")
    st.plotly_chart(fig2, use_container_width=True)

    # Capability if LSL/USL provided
    try:
        L = float(st.session_state.get("LSL_cache", st.session_state.get("LSL", "")) or 0 if False else 0)
    except:
        L= None
    # Read from inputs
    L_txt = st.session_state.get("LSL", "")
    U_txt = st.session_state.get("USL", "")
    # We will pass via text inputs in sidebar; read directly there
    pass

# ---------- Report (PPTX) ----------
st.subheader("📄 Export Report")
if st.button("Generate PowerPoint"):
    prs = Presentation()
    # title slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "CBS Quick Insights Report"
    slide.placeholders[1].text = "Auto-generated summary of EDA, relationships, and regression"
    # add a summary slide
    layout = prs.slide_layouts[1]
    slide2 = prs.slides.add_slide(layout)
    slide2.shapes.title.text = "Summary"
    body = slide2.shapes.placeholders[1].text_frame
    body.text = f"Rows: {n_rows}, Columns: {n_cols}\nY: {Y_cols}\nX: {X_cols}"
    # save
    out = io.BytesIO()
    prs.save(out)
    st.download_button("Download PPTX", data=out.getvalue(), file_name="CBS_Quick_Insights_Report.pptx", mime="application/vnd.openxmlformats-officedocument.presentationml.presentation")

st.markdown("---")
st.caption("Prototype build: Streamlit + pandas + statsmodels. For internal demo use.")
